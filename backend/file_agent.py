from __future__ import annotations

import csv
import json
import mimetypes
import os
from pathlib import Path
from typing import Any, Dict, List

try:
    from docx import Document
except Exception:
    Document = None

try:
    from openpyxl import load_workbook
except Exception:
    load_workbook = None

try:
    from pptx import Presentation
except Exception:
    Presentation = None

try:
    from pypdf import PdfReader
except Exception:
    PdfReader = None


TEXT_EXTENSIONS = {
    '.txt', '.md', '.json', '.jsonl', '.csv', '.py', '.js', '.jsx', '.ts', '.tsx', '.html', '.css', '.xml', '.yml', '.yaml', '.ini', '.log'
}


class FileAgent:
    def __init__(self, workspace_root: str | Path):
        self.workspace_root = Path(workspace_root)
        self.projects_dir = self.workspace_root / 'projects'

    def _safe_resolve(self, path: str | None) -> Path | None:
        if not path:
            return None
        p = Path(path).expanduser()
        if not p.is_absolute():
            p = (self.workspace_root / p).resolve()
        return p

    def detect_type(self, path: str) -> Dict[str, Any]:
        file_path = self._safe_resolve(path)
        if not file_path or not file_path.exists():
            return {"ok": False, "result": "Archivo no existe."}
        mime, _ = mimetypes.guess_type(str(file_path))
        return {
            "ok": True,
            "path": str(file_path),
            "name": file_path.name,
            "extension": file_path.suffix.lower(),
            "size_bytes": file_path.stat().st_size,
            "mime": mime or 'application/octet-stream',
            "is_text": file_path.suffix.lower() in TEXT_EXTENSIONS,
        }

    def read_file(self, path: str, max_chars: int = 50000) -> Dict[str, Any]:
        info = self.detect_type(path)
        if not info.get('ok'):
            return info
        file_path = Path(info['path'])
        ext = info.get('extension', '')
        try:
            if ext in {'.json', '.jsonl'}:
                content = file_path.read_text(encoding='utf-8', errors='ignore')[:max_chars]
                return {"ok": True, "path": str(file_path), "content": content, "summary": f"JSON leído: {file_path.name}"}
            if ext == '.pdf':
                if PdfReader is None:
                    return {"ok": False, "result": "Lectura PDF no disponible todavía: falta instalar pypdf en este entorno."}
                reader = PdfReader(str(file_path))
                parts = []
                for page in reader.pages[:20]:
                    parts.append((page.extract_text() or '').strip())
                content = '\n\n'.join([p for p in parts if p])[:max_chars]
                return {"ok": True, "path": str(file_path), "content": content, "summary": f"PDF leído con {len(reader.pages)} páginas."}
            if ext == '.docx':
                if Document is None:
                    return {"ok": False, "result": "Lectura DOCX no disponible todavía: falta instalar python-docx en este entorno."}
                doc = Document(str(file_path))
                content = '\n'.join([p.text for p in doc.paragraphs if p.text.strip()])[:max_chars]
                return {"ok": True, "path": str(file_path), "content": content, "summary": f"DOCX leído: {file_path.name}"}
            if ext == '.xlsx':
                if load_workbook is None:
                    return {"ok": False, "result": "Lectura XLSX no disponible todavía: falta instalar openpyxl en este entorno."}
                wb = load_workbook(str(file_path), read_only=True, data_only=True)
                lines = []
                for sheet in wb.worksheets[:5]:
                    lines.append(f"[Hoja] {sheet.title}")
                    for idx, row in enumerate(sheet.iter_rows(values_only=True)):
                        values = [str(v) for v in row if v is not None]
                        if values:
                            lines.append(' | '.join(values))
                        if idx >= 30:
                            break
                content = '\n'.join(lines)[:max_chars]
                return {"ok": True, "path": str(file_path), "content": content, "summary": f"XLSX leído con {len(wb.worksheets)} hojas."}
            if ext == '.pptx':
                if Presentation is None:
                    return {"ok": False, "result": "Lectura PPTX no disponible todavía: falta instalar python-pptx en este entorno."}
                prs = Presentation(str(file_path))
                lines = []
                for slide_index, slide in enumerate(prs.slides[:20], start=1):
                    lines.append(f"[Diapositiva {slide_index}]")
                    texts = []
                    for shape in slide.shapes:
                        if hasattr(shape, 'text') and str(shape.text).strip():
                            texts.append(str(shape.text).strip())
                    lines.extend(texts)
                content = '\n'.join(lines)[:max_chars]
                return {"ok": True, "path": str(file_path), "content": content, "summary": f"PPTX leído con {len(prs.slides)} diapositivas."}
            if ext == '.csv':
                rows = []
                with open(file_path, 'r', encoding='utf-8', errors='ignore', newline='') as f:
                    reader = csv.reader(f)
                    for idx, row in enumerate(reader):
                        rows.append(row)
                        if idx >= 50:
                            break
                preview = '\n'.join([', '.join(r) for r in rows])[:max_chars]
                return {"ok": True, "path": str(file_path), "content": preview, "summary": f"CSV leído con {len(rows)} filas de vista previa."}
            if ext in TEXT_EXTENSIONS:
                content = file_path.read_text(encoding='utf-8', errors='ignore')[:max_chars]
                return {"ok": True, "path": str(file_path), "content": content, "summary": f"Texto leído: {file_path.name}"}
            return {"ok": False, "result": f"Tipo de archivo no soportado aún para lectura directa: {ext or 'sin extensión'}"}
        except Exception as e:
            return {"ok": False, "result": f"Error leyendo archivo: {e}"}

    def list_files(self, base_path: str | None = None, limit: int = 100) -> Dict[str, Any]:
        target = self._safe_resolve(base_path) if base_path else self.projects_dir
        if not target or not target.exists():
            return {"ok": False, "result": "Ruta no existe."}
        items: List[Dict[str, Any]] = []
        try:
            if target.is_file():
                target = target.parent
            for path in sorted(target.rglob('*')):
                if len(items) >= limit:
                    break
                if path.is_file():
                    items.append({
                        'name': path.name,
                        'path': str(path),
                        'extension': path.suffix.lower(),
                        'size_bytes': path.stat().st_size,
                    })
            return {"ok": True, "items": items, "count": len(items), "result": f"{len(items)} archivos listados."}
        except Exception as e:
            return {"ok": False, "result": f"Error listando archivos: {e}"}

    def recent_project_files(self, limit: int = 20) -> Dict[str, Any]:
        items = []
        try:
            for path in self.projects_dir.rglob('*'):
                if path.is_file():
                    items.append(path)
            items = sorted(items, key=lambda p: p.stat().st_mtime, reverse=True)[:limit]
            return {
                "ok": True,
                "items": [
                    {
                        'name': p.name,
                        'path': str(p),
                        'extension': p.suffix.lower(),
                        'size_bytes': p.stat().st_size,
                    }
                    for p in items
                ],
                "count": len(items),
            }
        except Exception as e:
            return {"ok": False, "result": f"Error obteniendo archivos recientes: {e}"}

    def suggest_recent_attachments(self, limit: int = 8, allowed_extensions: List[str] | None = None) -> Dict[str, Any]:
        recent = self.recent_project_files(max(limit * 3, limit))
        if not recent.get('ok'):
            return recent
        allowed = {ext.lower() for ext in (allowed_extensions or ['.pdf', '.docx', '.xlsx', '.pptx', '.png', '.jpg', '.jpeg', '.txt', '.md', '.csv', '.stl'])}
        items = [item for item in recent.get('items', []) if str(item.get('extension', '')).lower() in allowed][:limit]
        return {
            'ok': True,
            'items': items,
            'count': len(items),
            'result': f"{len(items)} adjuntos recientes sugeridos.",
        }

    def prepare_email_attachments(self, paths: List[str], max_size_mb: int = 20) -> Dict[str, Any]:
        resolved = []
        errors = []
        max_bytes = max_size_mb * 1024 * 1024
        total = 0
        for raw in paths or []:
            file_path = self._safe_resolve(raw)
            if not file_path or not file_path.exists() or not file_path.is_file():
                errors.append(f"No existe: {raw}")
                continue
            size = file_path.stat().st_size
            if size > max_bytes:
                errors.append(f"Muy grande: {file_path.name} ({round(size / (1024*1024), 2)} MB)")
                continue
            total += size
            if total > max_bytes:
                errors.append(f"Límite total excedido al agregar {file_path.name}")
                break
            mime, _ = mimetypes.guess_type(str(file_path))
            resolved.append({
                'path': str(file_path),
                'name': file_path.name,
                'size_bytes': size,
                'mime': mime or 'application/octet-stream',
            })
        return {
            'ok': len(resolved) > 0,
            'attachments': resolved,
            'errors': errors,
            'count': len(resolved),
            'total_size_bytes': total,
            'result': f"Adjuntos preparados: {len(resolved)}",
        }
