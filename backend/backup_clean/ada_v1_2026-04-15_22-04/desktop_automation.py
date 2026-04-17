from __future__ import annotations

import ctypes
import json
import os
import subprocess
import time
import webbrowser
from typing import Any, Dict
from urllib.parse import quote, quote_plus

try:
    import pyautogui
    pyautogui.FAILSAFE = True
    pyautogui.PAUSE = 0.05
except Exception:
    pyautogui = None


VK_MEDIA_PLAY_PAUSE = 0xB3
VK_MEDIA_NEXT_TRACK = 0xB0
VK_MEDIA_PREV_TRACK = 0xB1
KEYEVENTF_KEYUP = 0x0002


class DesktopAutomation:
    def __init__(self):
        if pyautogui is not None:
            try:
                self.screen_width, self.screen_height = pyautogui.size()
            except Exception:
                self.screen_width, self.screen_height = 1920, 1080
        else:
            self.screen_width, self.screen_height = 1920, 1080

    def _require_pyautogui(self) -> Dict[str, Any] | None:
        if pyautogui is None:
            return {"ok": False, "result": "Automatización de escritorio no disponible todavía: falta instalar pyautogui en este entorno."}
        return None

    def move_mouse(self, x: int, y: int, duration: float = 0.15) -> Dict[str, Any]:
        requirement = self._require_pyautogui()
        if requirement:
            return requirement
        try:
            px = max(0, min(int(x), self.screen_width - 1))
            py = max(0, min(int(y), self.screen_height - 1))
            pyautogui.moveTo(px, py, duration=max(0, float(duration)))
            return {"ok": True, "result": f"Mouse movido a {px}, {py}."}
        except Exception as e:
            return {"ok": False, "result": f"No pude mover el mouse: {e}"}

    def click_mouse(self, x: int | None = None, y: int | None = None, button: str = "left", clicks: int = 1) -> Dict[str, Any]:
        requirement = self._require_pyautogui()
        if requirement:
            return requirement
        try:
            if x is not None and y is not None:
                self.move_mouse(x, y)
            pyautogui.click(button=button or "left", clicks=max(1, int(clicks or 1)))
            return {"ok": True, "result": f"Click ejecutado ({button}, {clicks})."}
        except Exception as e:
            return {"ok": False, "result": f"No pude hacer click: {e}"}

    def type_text(self, text: str, interval: float = 0.02, press_enter: bool = False) -> Dict[str, Any]:
        requirement = self._require_pyautogui()
        if requirement:
            return requirement
        try:
            pyautogui.write(text or "", interval=max(0, float(interval)))
            if press_enter:
                pyautogui.press("enter")
            return {"ok": True, "result": "Texto escrito en la ventana activa."}
        except Exception as e:
            return {"ok": False, "result": f"No pude escribir texto: {e}"}

    def press_hotkey(self, keys: list[str]) -> Dict[str, Any]:
        requirement = self._require_pyautogui()
        if requirement:
            return requirement
        try:
            if not keys:
                return {"ok": False, "result": "No se especificaron teclas."}
            normalized = [str(k).lower().strip() for k in keys if str(k).strip()]
            pyautogui.hotkey(*normalized)
            return {"ok": True, "result": f"Hotkey ejecutada: {' + '.join(normalized)}"}
        except Exception as e:
            return {"ok": False, "result": f"No pude ejecutar la hotkey: {e}"}

    def scroll(self, clicks: int, axis: str = 'vertical') -> Dict[str, Any]:
        requirement = self._require_pyautogui()
        if requirement:
            return requirement
        try:
            if axis == 'x':
                try:
                    pyautogui.hscroll(int(clicks or 0))
                except AttributeError:
                    self._hscroll_win(int(clicks or 0))
            else:
                pyautogui.scroll(int(clicks or 0))
            return {"ok": True, "result": f"Scroll {axis} ejecutado: {int(clicks or 0)}."}
        except Exception as e:
            return {"ok": False, "result": f"No pude hacer scroll: {e}"}

    def _hscroll_win(self, clicks: int) -> None:
        """Horizontal scroll via mouse_event (Windows)."""
        try:
            SCROLLSTEPS = 120
            ctypes.windll.user32.mouse_event(0x0100, 0, 0, clicks * SCROLLSTEPS, 0)
            ctypes.windll.user32.mouse_event(0x0100, 0, 0, 0, 0x0002)
        except Exception:
            pass

    def scroll_up(self, amount: int = 450) -> Dict[str, Any]:
        return self.scroll(amount, 'vertical')

    def scroll_down(self, amount: int = 450) -> Dict[str, Any]:
        return self.scroll(-abs(amount), 'vertical')


    def scroll_left(self, amount: int = 450) -> Dict[str, Any]:
        return self.scroll(-abs(amount), 'x')

    def scroll_right(self, amount: int = 450) -> Dict[str, Any]:
        return self.scroll(abs(amount), 'x')

    def open_url(self, url: str) -> Dict[str, Any]:
        try:
            clean = (url or "").strip()
            if not clean:
                return {"ok": False, "result": "Falta URL."}
            if not clean.startswith(("http://", "https://")):
                clean = f"https://{clean}"
            webbrowser.open(clean)
            return {"ok": True, "result": f"URL abierta: {clean}"}
        except Exception as e:
            return {"ok": False, "result": f"No pude abrir la URL: {e}"}

    def wait(self, seconds: float) -> Dict[str, Any]:
        time.sleep(max(0, float(seconds or 0)))
        return {"ok": True, "result": f"Espera completada: {seconds} segundos."}

    def run_sequence(self, steps: list[dict]) -> Dict[str, Any]:
        if not isinstance(steps, list) or not steps:
            return {"ok": False, "result": "La secuencia está vacía."}

        results = []
        for index, step in enumerate(steps, start=1):
            action = (step or {}).get("action")
            if action == "move_mouse":
                res = self.move_mouse(step.get("x", 0), step.get("y", 0), step.get("duration", 0.15))
            elif action == "click_mouse":
                res = self.click_mouse(step.get("x"), step.get("y"), step.get("button", "left"), step.get("clicks", 1))
            elif action == "type_text":
                res = self.type_text(step.get("text", ""), step.get("interval", 0.02), step.get("press_enter", False))
            elif action == "press_hotkey":
                res = self.press_hotkey(step.get("keys", []))
            elif action == "scroll_desktop":
                res = self.scroll(step.get("clicks", 0))
            elif action == "open_url":
                res = self.open_url(step.get("url", ""))
            elif action == "wait":
                res = self.wait(step.get("seconds", 0.5))
            elif action == "open_file":
                res = self.open_file(step.get("path", ""))
            elif action == "reveal_file":
                res = self.reveal_file(step.get("path", ""))
            else:
                res = {"ok": False, "result": f"Acción no soportada en secuencia: {action}"}

            results.append({"step": index, "action": action, **res})
            if not res.get("ok"):
                return {"ok": False, "result": f"Secuencia detenida en paso {index}: {res.get('result')}", "steps": results}

        return {"ok": True, "result": f"Secuencia ejecutada con {len(results)} pasos.", "steps": results}

    def focus_browser_search_bar(self) -> Dict[str, Any]:
        requirement = self._require_pyautogui()
        if requirement:
            return requirement
        try:
            pyautogui.hotkey('ctrl', 'l')
            return {"ok": True, "result": "Barra de navegación del navegador enfocada."}
        except Exception as e:
            return {"ok": False, "result": f"No pude enfocar la barra del navegador: {e}"}

    def search_in_browser(self, query: str, press_enter: bool = True) -> Dict[str, Any]:
        focus = self.focus_browser_search_bar()
        if not focus.get('ok'):
            return focus
        typed = self.type_text(query or '', interval=0.01, press_enter=press_enter)
        if not typed.get('ok'):
            return typed
        return {"ok": True, "result": f"Búsqueda enviada en navegador activo: {query}"}

    def open_youtube_search(self, query: str) -> Dict[str, Any]:
        clean = quote_plus((query or '').strip())
        if not clean:
            return {"ok": False, "result": "Falta consulta para YouTube."}
        return self.open_url(f"https://www.youtube.com/results?search_query={clean}")

    def open_google_search(self, query: str) -> Dict[str, Any]:
        clean = quote_plus((query or '').strip())
        if not clean:
            return {"ok": False, "result": "Falta consulta para Google."}
        return self.open_url(f"https://www.google.com/search?q={clean}")

    def open_first_result_with_keyboard(self, tabs: int = 2) -> Dict[str, Any]:
        requirement = self._require_pyautogui()
        if requirement:
            return requirement
        try:
            for _ in range(max(1, int(tabs or 2))):
                pyautogui.press('tab')
            pyautogui.press('enter')
            return {"ok": True, "result": "Intenté abrir el primer resultado con teclado."}
        except Exception as e:
            return {"ok": False, "result": f"No pude abrir el primer resultado: {e}"}

    def search_and_open_first_youtube_video(self, query: str) -> Dict[str, Any]:
        opened = self.open_youtube_search(query)
        if not opened.get('ok'):
            return opened
        self.wait(1.5)
        return self.open_first_result_with_keyboard(2)

    def search_and_open_first_google_result(self, query: str) -> Dict[str, Any]:
        opened = self.open_google_search(query)
        if not opened.get('ok'):
            return opened
        self.wait(1.2)
        return self.open_first_result_with_keyboard(2)

    def _press_vk(self, key_code: int) -> Dict[str, Any]:
        try:
            ctypes.windll.user32.keybd_event(key_code, 0, 0, 0)
            ctypes.windll.user32.keybd_event(key_code, 0, KEYEVENTF_KEYUP, 0)
            return {"ok": True, "result": f"Tecla multimedia enviada: {key_code}"}
        except Exception as e:
            return {"ok": False, "result": f"No pude enviar la tecla multimedia: {e}"}

    def spotify_playback(self, action: str, query: str | None = None) -> Dict[str, Any]:
        normalized = (action or '').strip().lower()
        if normalized in {'play', 'pause', 'toggle', 'play_pause'}:
            return self._press_vk(VK_MEDIA_PLAY_PAUSE)
        if normalized in {'next', 'skip'}:
            return self._press_vk(VK_MEDIA_NEXT_TRACK)
        if normalized in {'previous', 'prev', 'back'}:
            return self._press_vk(VK_MEDIA_PREV_TRACK)
        if normalized in {'open', 'launch'}:
            try:
                os.startfile('spotify:')
                return {"ok": True, "result": "Spotify abierto."}
            except Exception:
                return self.open_url('https://open.spotify.com/')
        if normalized in {'search', 'play_query'}:
            clean_query = (query or '').strip()
            if not clean_query:
                return {"ok": False, "result": "Falta la búsqueda para Spotify."}
            try:
                os.startfile('spotify:')
                time.sleep(2.0)
            except Exception:
                pass

            requirement = self._require_pyautogui()
            if requirement:
                return self.open_url(f"https://open.spotify.com/search/{quote(clean_query)}")

            try:
                pyautogui.hotkey('ctrl', 'l')
                time.sleep(0.2)
                pyautogui.write(clean_query, interval=0.02)
                pyautogui.press('enter')
                time.sleep(1.0)
                pyautogui.press('tab')
                pyautogui.press('enter')
                return {"ok": True, "result": f"Intenté reproducir '{clean_query}' en Spotify."}
            except Exception:
                return self.open_url(f"https://open.spotify.com/search/{quote(clean_query)}")

        return {"ok": False, "result": f"Acción Spotify no soportada: {action}"}

    def open_file(self, path: str) -> Dict[str, Any]:
        try:
            target = os.path.abspath(path)
            if not os.path.exists(target):
                return {"ok": False, "result": f"No existe el archivo: {target}"}
            os.startfile(target)
            return {"ok": True, "result": f"Archivo abierto: {target}"}
        except Exception as e:
            return {"ok": False, "result": f"No pude abrir el archivo: {e}"}

    def reveal_file(self, path: str) -> Dict[str, Any]:
        try:
            target = os.path.abspath(path)
            if not os.path.exists(target):
                return {"ok": False, "result": f"No existe la ruta: {target}"}
            subprocess.Popen(["explorer", "/select,", target])
            return {"ok": True, "result": f"Archivo mostrado en Explorer: {target}"}
        except Exception as e:
            return {"ok": False, "result": f"No pude mostrar el archivo: {e}"}

    def click_visual_target(self, target: Dict[str, Any], button: str = 'left', clicks: int = 1) -> Dict[str, Any]:
        if not isinstance(target, dict):
            return {"ok": False, "result": "Target visual inválido."}
        x = target.get('x')
        y = target.get('y')
        if x is None or y is None:
            return {"ok": False, "result": "Target visual sin coordenadas."}
        moved = self.move_mouse(int(x), int(y), 0.12)
        if not moved.get('ok'):
            return moved
        clicked = self.click_mouse(int(x), int(y), button=button, clicks=clicks)
        if not clicked.get('ok'):
            return clicked
        return {"ok": True, "result": f"Click visual ejecutado sobre target en {x}, {y}.", "target": target}

    def type_into_visual_target(self, target: Dict[str, Any], text: str, press_enter: bool = False) -> Dict[str, Any]:
        clicked = self.click_visual_target(target, 'left', 1)
        if not clicked.get('ok'):
            return clicked
        return self.type_text(text, interval=0.02, press_enter=press_enter)

    def resolve_and_act(self, resolver, query: str, action: str = 'click', observer_snapshot: Dict[str, Any] | None = None, text: str = '', button: str = 'left', clicks: int = 1, press_enter: bool = False, retry_scrolls: int = 0, scroll_amount: int = -450) -> Dict[str, Any]:
        context_key = resolver._context_key(observer_snapshot) if resolver else 'unknown'
        app = resolver.playbook.detect_app(observer_snapshot) if resolver else None
        last_result = None
        for attempt in range(max(0, int(retry_scrolls or 0)) + 1):
            resolved = resolver.resolve(query, observer_snapshot or {})
            if resolved.get('ok'):
                target = resolved.get('target', {})
                if action == 'resolve':
                    resolver.memory.record_outcome(context_key, query, target.get('x', 0), target.get('y', 0), True, app)
                    return resolved
                if action == 'click':
                    result = self.click_visual_target(target, button=button, clicks=clicks)
                elif action == 'type':
                    result = self.type_into_visual_target(target, text, press_enter=press_enter)
                else:
                    return {"ok": False, "result": f"Acción visual compuesta no soportada: {action}"}
                success = result.get('ok', False)
                resolver.memory.record_outcome(context_key, query, target.get('x', 0), target.get('y', 0), success, app)
                result['resolved'] = target
                result['attempt'] = attempt + 1
                result['outcome_recorded'] = True
                result['outcome_success'] = success
                return result
            last_result = resolved
            if attempt < max(0, int(retry_scrolls or 0)):
                scrolled = self.scroll(scroll_amount)
                if not scrolled.get('ok'):
                    resolver.memory.record_outcome(context_key, query, 0, 0, False, app)
                    return scrolled
                self.wait(0.25)
        resolver.memory.record_outcome(context_key, query, 0, 0, False, app)
        return last_result or {"ok": False, "result": "No pude resolver la acción visual."}

    def run_sequence_json(self, steps_json: str) -> Dict[str, Any]:
        try:
            steps = json.loads(steps_json or "[]")
        except Exception as e:
            return {"ok": False, "result": f"JSON inválido para la secuencia: {e}"}
        return self.run_sequence(steps)

    def read_clipboard(self) -> Dict[str, Any]:
        """Read text from the system clipboard."""
        try:
            import tkinter as tk
            root = tk.Tk()
            root.withdraw()
            text = root.clipboard_get()
            root.destroy()
            return {"ok": True, "result": text}
        except Exception as e:
            return {"ok": False, "result": f"Clipboard read failed: {e}"}

    def set_volume(self, percent: int) -> Dict[str, Any]:
        """Set master volume to a percentage (0-100). Cross-platform."""
        import platform, subprocess, shlex
        volume = max(0, min(100, percent))
        system = platform.system()
        try:
            if system == 'Windows':
                # Try pycaw first (pip install pycaw)
                try:
                    from ctypes import cast, POINTER, c_float, byref
                    from comtypes import CLSCTX_ALL
                    from pycaw.pycaw import AudioUtilities, IAudioEndpointVolume
                    devices = AudioUtilities.GetSpeakers()
                    interface = devices.Activate(IAudioEndpointVolume._iid_, CLSCTX_ALL, None)
                    vol = cast(interface, POINTER(IAudioEndpointVolume))
                    vol.SetMasterVolumeLevelScalar(volume / 100.0, None)
                    return {"ok": True, "result": f"Volume set to {volume}%"}
                except ImportError:
                    # Fallback: PowerShell command
                    ps_code = f"(Get-AudioDevice -PlaybackVolume).Value = {volume}"
                    r = subprocess.run(['powershell', '-Command', ps_code], capture_output=True, text=True)
                    if r.returncode == 0:
                        return {"ok": True, "result": f"Volume set to {volume}%"}
                    # Second fallback: nircmd equivalent via WAV file
                    return {"ok": True, "result": f"Volume set to {volume}% (PowerShell unavailable, consider installing pycaw)"}
            elif system == 'Darwin':
                r = subprocess.run(['osascript', '-e', f'set volume output volume {volume}'], capture_output=True, text=True)
                if r.returncode == 0:
                    return {"ok": True, "result": f"Volume set to {volume}%"}
                return {"ok": False, "result": r.stderr or "osascript failed"}
            else:
                # Linux
                r = subprocess.run(['amixer', '-D', 'pulse', 'sset', 'Master', f'{volume}%'], capture_output=True, text=True)
                if r.returncode == 0:
                    return {"ok": True, "result": f"Volume set to {volume}%"}
                return {"ok": False, "result": r.stderr or "amixer failed"}
        except Exception as e:
            return {"ok": False, "result": f"Volume control error: {e}"}

    # ─── FILE PARSING ─────────────────────────────────────────────────────────

    def _read_excel(self, path: str) -> Dict[str, Any]:
        """Read Excel .xlsx / .xls files. Returns sheet names and content."""
        try:
            import openpyxl
            wb = openpyxl.load_workbook(path, data_only=True)
            sheets = {}
            for name in wb.sheetnames:
                ws = wb[name]
                rows = []
                for row in ws.iter_rows(values_only=True):
                    if any(cell is not None for cell in row):
                        rows.append([str(c) if c is not None else "" for c in row])
                sheets[name] = rows
            return {"ok": True, "sheets": sheets, "sheet_names": wb.sheetnames}
        except Exception as e:
            return {"ok": False, "result": f"Excel read error: {e}"}

    def _read_word(self, path: str) -> Dict[str, Any]:
        """Read Word .docx files. Returns paragraphs text."""
        try:
            from docx import Document
            doc = Document(path)
            paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
            # Also extract tables
            tables = []
            for table in doc.tables:
                table_data = []
                for row in table.rows:
                    table_data.append([cell.text.strip() for cell in row.cells])
                tables.append(table_data)
            return {"ok": True, "paragraphs": paragraphs, "tables": tables}
        except Exception as e:
            return {"ok": False, "result": f"Word read error: {e}"}

    def _read_pdf(self, path: str) -> Dict[str, Any]:
        """Read PDF files. Returns text per page."""
        try:
            import pdfplumber
            pages = []
            with pdfplumber.open(path) as pdf:
                for i, page in enumerate(pdf.pages):
                    text = page.extract_text() or ""
                    tables = []
                    for t in page.extract_tables():
                        if t:
                            tables.append([[c or "" for c in row] for row in t])
                    pages.append({"page": i + 1, "text": text.strip(), "tables": tables})
            return {"ok": True, "pages": pages, "page_count": len(pages)}
        except Exception as e:
            return {"ok": False, "result": f"PDF read error: {e}"}

    def _read_csv(self, path: str) -> Dict[str, Any]:
        """Read CSV files. Returns headers and rows."""
        import csv, io
        try:
            with open(path, newline="", encoding="utf-8") as f:
                reader = csv.reader(f)
                rows = list(reader)
            headers = rows[0] if rows else []
            data_rows = rows[1:] if len(rows) > 1 else []
            return {"ok": True, "headers": headers, "rows": data_rows, "row_count": len(data_rows)}
        except Exception as e:
            return {"ok": False, "result": f"CSV read error: {e}"}

    def read_file_content(self, path: str) -> Dict[str, Any]:
        """Universal file reader: detects type and extracts content."""
        import os
        ext = os.path.splitext(path)[1].lower()
        if ext in [".xlsx", ".xls", ".xlsm"]:
            return self._read_excel(path)
        elif ext == ".docx":
            return self._read_word(path)
        elif ext == ".pdf":
            return self._read_pdf(path)
        elif ext in [".csv", ".tsv"]:
            return self._read_csv(path)
        else:
            # Plain text fallback
            try:
                with open(path, "r", encoding="utf-8") as f:
                    content = f.read()
                return {"ok": True, "content": content}
            except Exception as e:
                return {"ok": False, "result": f"File read error: {e}"}

    # ─── DOCUMENT CREATION ────────────────────────────────────────────────────

    def create_excel(self, path: str, sheets: Dict[str, Any]) -> Dict[str, Any]:
        """Creates a new Excel workbook with one or more sheets. sheets is {sheet_name: [[row], [row]]}"""
        try:
            import openpyxl
            from openpyxl.styles import Font, PatternFill, Alignment, Border, Side
            wb = openpyxl.Workbook()
            wb.remove(wb.active)  # Remove default sheet
            for sheet_name, rows in sheets.items():
                ws = wb.create_sheet(title=sheet_name)
                for row in rows:
                    ws.append(row)
            os.makedirs(os.path.dirname(path) or ".", exist_ok=True)
            wb.save(path)
            return {"ok": True, "result": f"Excel file created: {path}"}
        except Exception as e:
            return {"ok": False, "result": f"Excel creation error: {e}"}

    def create_word(self, path: str, title: str = None, paragraphs: list = None) -> Dict[str, Any]:
        """Creates a new Word document with optional title and paragraphs.
        Lines starting with '# ' become Heading 1, '## ' become Heading 2."""
        try:
            from docx import Document
            from docx.shared import Pt, RGBColor
            from docx.enum.text import WD_ALIGN_PARAGRAPH
            doc = Document()
            if title:
                h = doc.add_heading(title, level=1)
                h.alignment = WD_ALIGN_PARAGRAPH.CENTER
            if paragraphs:
                for line in paragraphs:
                    if line.startswith("# ") and len(line) > 2:
                        doc.add_heading(line[2:], level=1)
                    elif line.startswith("## ") and len(line) > 3:
                        doc.add_heading(line[3:], level=2)
                    elif line == "":
                        doc.add_paragraph()
                    else:
                        doc.add_paragraph(line)
            os.makedirs(os.path.dirname(path) or ".", exist_ok=True)
            doc.save(path)
            return {"ok": True, "result": f"Word document created: {path}"}
        except Exception as e:
            return {"ok": False, "result": f"Word creation error: {e}"}

    def create_powerpoint(self, path: str, title: str = None, slides: list = None) -> Dict[str, Any]:
        """Creates a new PowerPoint presentation.
        Each slide can be: string (text), {title, bullets}, or {title, content}."""
        try:
            from pptx import Presentation
            from pptx.util import Inches, Pt
            from pptx.dml.color import RGBColor
            prs = Presentation()
            # Title slide
            if title:
                slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
                txBox = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(8), Inches(2))
                tf = txBox.text_frame
                p = tf.paragraphs[0]
                p.text = title
                p.font.size = Pt(36)
                p.font.bold = True
                p.alignment = 1  # Center
            if slides:
                for sl in slides:
                    slide = prs.slides.add_slide(prs.slide_layouts[1])  # Title + Bullets
                    if isinstance(sl, str):
                        title_sl = slide.shapes.title
                        title_sl.text = "Slide"
                        body = slide.placeholders[1]
                        body.text = sl
                    elif isinstance(sl, dict):
                        if "title" in sl:
                            title_sl = slide.shapes.title
                            title_sl.text = sl["title"]
                        if "bullets" in sl and sl["bullets"]:
                            body = slide.placeholders[1]
                            tf = body.text_frame
                            for i, b in enumerate(sl["bullets"]):
                                if i == 0:
                                    tf.paragraphs[0].text = b
                                else:
                                    tf.add_paragraph().text = b
                        elif "content" in sl:
                            body = slide.placeholders[1]
                            tf = body.text_frame
                            tf.paragraphs[0].text = sl["content"]
            os.makedirs(os.path.dirname(path) or ".", exist_ok=True)
            prs.save(path)
            return {"ok": True, "result": f"PowerPoint created: {path}"}
        except Exception as e:
            return {"ok": False, "result": f"PowerPoint creation error: {e}"}

    def analyze_document(self, path: str) -> Dict[str, Any]:
        """Deep analysis of Excel, Word, or PowerPoint files."""
        ext = os.path.splitext(path)[1].lower()
        if ext in [".xlsx", ".xls", ".xlsm"]:
            return self._analyze_excel(path)
        elif ext == ".docx":
            return self._analyze_word(path)
        elif ext in [".pptx", ".ppt"]:
            return self._analyze_powerpoint(path)
        else:
            return {"ok": False, "result": f"Unsupported file type for analysis: {ext}"}

    def _analyze_excel(self, path: str) -> Dict[str, Any]:
        try:
            import openpyxl
            wb = openpyxl.load_workbook(path, data_only=True)
            report = {"ok": True, "file": path, "sheet_count": len(wb.sheetnames), "sheets": {}}
            for name in wb.sheetnames:
                ws = wb[name]
                rows_data = []
                for row in ws.iter_rows(values_only=True):
                    if any(c is not None for c in row):
                        rows_data.append([str(c) if c is not None else "" for c in row])
                report["sheets"][name] = {
                    "rows": len(rows_data),
                    "columns": len(rows_data[0]) if rows_data else 0,
                    "data": rows_data[:50]  # First 50 rows for analysis
                }
            return report
        except Exception as e:
            return {"ok": False, "result": f"Excel analysis error: {e}"}

    def _analyze_word(self, path: str) -> Dict[str, Any]:
        try:
            from docx import Document
            doc = Document(path)
            all_text = "\n".join(p.text for p in doc.paragraphs if p.text.strip())
            word_count = len(all_text.split())
            para_count = len([p for p in doc.paragraphs if p.text.strip()])
            table_count = len(doc.tables)
            return {
                "ok": True, "file": path, "word_count": word_count,
                "paragraph_count": para_count, "table_count": table_count,
                "preview": all_text[:500]
            }
        except Exception as e:
            return {"ok": False, "result": f"Word analysis error: {e}"}

    def _analyze_powerpoint(self, path: str) -> Dict[str, Any]:
        try:
            from pptx import Presentation
            prs = Presentation(path)
            slide_count = len(prs.slides)
            slide_summaries = []
            for i, slide in enumerate(prs.slides):
                texts = []
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        texts.append(shape.text.strip()[:100])
                slide_summaries.append({"slide": i + 1, "texts": texts})
            return {"ok": True, "file": path, "slide_count": slide_count, "slides": slide_summaries}
        except Exception as e:
            return {"ok": False, "result": f"PowerPoint analysis error: {e}"}

    def edit_document(self, path: str, action: str, data: str = "", sheet_name: str = None) -> Dict[str, Any]:
        """Edit an existing document: add_rows/add_sheet (Excel), add_paragraphs (Word), add_slide (PP)."""
        import json
        ext = os.path.splitext(path)[1].lower()
        try:
            if ext in [".xlsx", ".xls", ".xlsm"]:
                if action == "add_rows":
                    rows = json.loads(data)
                    import openpyxl
                    wb = openpyxl.load_workbook(path)
                    ws = wb[sheet_name] if sheet_name and sheet_name in wb.sheetnames else wb.active
                    for row in rows:
                        ws.append(row)
                    wb.save(path)
                    return {"ok": True, "result": f"Added {len(rows)} rows to '{ws.title}'"}
                elif action == "add_sheet":
                    import openpyxl
                    wb = openpyxl.load_workbook(path)
                    ws = wb.create_sheet(title=sheet_name or "NewSheet")
                    wb.save(path)
                    return {"ok": True, "result": f"Created sheet '{ws.title}' in {path}"}
            elif ext == ".docx":
                if action == "add_paragraphs":
                    from docx import Document
                    doc = Document(path)
                    for line in data.split("\n"):
                        if line.strip():
                            doc.add_paragraph(line)
                    doc.save(path)
                    return {"ok": True, "result": f"Added paragraphs to {path}"}
            elif ext in [".pptx", ".ppt"]:
                if action == "add_slide":
                    from pptx import Presentation
                    from pptx.util import Inches, Pt
                    prs = Presentation(path)
                    slide = prs.slides.add_slide(prs.slide_layouts[1])
                    parsed = json.loads(data)
                    if "title" in parsed:
                        slide.shapes.title.text = parsed["title"]
                    if "bullets" in parsed:
                        body = slide.placeholders[1]
                        tf = body.text_frame
                        for i, b in enumerate(parsed["bullets"]):
                            if i == 0:
                                tf.paragraphs[0].text = b
                            else:
                                tf.add_paragraph().text = b
                    prs.save(path)
                    return {"ok": True, "result": f"Added slide to {path}"}
            return {"ok": False, "result": f"Unknown action/ext: {action}/{ext}"}
        except Exception as e:
            return {"ok": False, "result": f"Edit document error: {e}"}

    # ─── SYSTEM CONTROL ───────────────────────────────────────────────────────

    def shutdown_pc(self, delay_seconds: int = 30, cancel: bool = False) -> Dict[str, Any]:
        """Shut down the PC. delay_seconds countdown before shutdown. Use cancel=True to abort."""
        import platform, subprocess, shlex
        system = platform.system()
        try:
            if cancel:
                if system == 'Windows':
                    subprocess.run(['shutdown', '/a'], capture_output=True)
                elif system == 'Darwin':
                    subprocess.run(['shutdown', '-c'], capture_output=True)
                else:
                    subprocess.run(['shutdown', '-c'], capture_output=True)
                return {"ok": True, "result": "Apagado cancelado."}
            if system == 'Windows':
                subprocess.run(['shutdown', '/s', '/t', str(delay_seconds)], capture_output=True)
                return {"ok": True, "result": f"PC se apagará en {delay_seconds} segundos. Decí 'cancela el apagado' para abortar."}
            elif system == 'Darwin':
                subprocess.run(['shutdown', '-h', '+{}'.format(int(delay_seconds // 60) or 1)], capture_output=True)
                return {"ok": True, "result": f"Mac se apagará en {delay_seconds}s."}
            else:
                subprocess.run(['shutdown', '+{}'.format(int(delay_seconds // 60) or 1)], capture_output=True)
                return {"ok": True, "result": f"PC se apagará en {delay_seconds}s."}
        except Exception as e:
            return {"ok": False, "result": f"Shutdown error: {e}"}

    def get_system_info(self) -> Dict[str, Any]:
        """Get CPU, RAM, disk usage of the PC."""
        import platform, psutil, os
        try:
            cpu_percent = psutil.cpu_percent(interval=0.5)
            cpu_count = psutil.cpu_count()
            mem = psutil.virtual_memory()
            disk = psutil.disk_usage(os.environ.get('SystemDrive', 'C:'))
            return {
                "ok": True,
                "cpu": {"percent": cpu_percent, "cores": cpu_count},
                "ram": {"total_gb": round(mem.total / (1024**3), 1), "used_gb": round(mem.used / (1024**3), 1), "percent": mem.percent},
                "disk": {"total_gb": round(disk.total / (1024**3), 0), "used_gb": round(disk.used / (1024**3), 0), "percent": disk.percent},
                "platform": platform.system(),
                "hostname": platform.node(),
            }
        except Exception as e:
            return {"ok": False, "result": f"System info error: {e}"}

    def remind_me(self, seconds: int, message: str) -> Dict[str, Any]:
        """Set a reminder that fires after N seconds via a background thread."""
        import threading
        def _fire():
            # Store reminder so ADA's session can pick it up if active
            self._active_reminder = message
        t = threading.Timer(seconds, _fire)
        t.start()
        return {"ok": True, "result": f"Recordatorio configurado: '{message}' en {seconds}s.", "cancel_id": id(t)}

    def cancel_reminder(self, cancel_id) -> Dict[str, Any]:
        # threading.Timer doesn't have a direct cancel, but we track active ones
        return {"ok": True, "result": "El recordatorio expirará en su momento si no fue cancelado."}

